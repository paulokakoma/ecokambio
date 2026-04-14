#!/usr/bin/env python3
"""
Gerador da Apresentação PowerPoint EcoKambio para Investidores 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE      = os.path.dirname(os.path.abspath(__file__))
SHOTS     = os.path.join(BASE, "screenshots")
OUT_FILE  = os.path.join(BASE, "EcoKambio_Investidores_2026.pptx")
LOGO_PATH = "/Users/av/Documents/Projetos/ecokambio-main/public/assets/main-logo.png"

IMG = {
    "mercado_informal":   os.path.join(SHOTS, "mercado_informal.png"),
    "homepage_produtos":  os.path.join(SHOTS, "homepage_produtos.png"),
    "calculadora":        os.path.join(SHOTS, "calculadora.png"),
    "parceiros_produtos": os.path.join(SHOTS, "parceiros_produtos.png"),
    "ecoflix":            os.path.join(SHOTS, "ecoflix.png"),
    "visa_hero":          os.path.join(SHOTS, "visa_hero.png"),
    "produto_detalhe":    os.path.join(SHOTS, "produto_detalhe.png"),
    # kept for metrics slide
    "faq":                os.path.join(SHOTS, "faq.png"),
}

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0D, 0x11, 0x17)   # quase preto
C_BG_CARD   = RGBColor(0x16, 0x1B, 0x22)   # cinza escuro
C_GREEN     = RGBColor(0x2A, 0x9D, 0x6F)   # verde eco
C_GREEN_L   = RGBColor(0x3B, 0xBD, 0x87)   # verde claro
C_GOLD      = RGBColor(0xF4, 0xA2, 0x27)   # dourado
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED     = RGBColor(0x8B, 0x94, 0x9E)
C_ACCENT    = RGBColor(0x1F, 0x29, 0x37)   # fundo de cards

# ── Slide size: 16:9 widescreen ───────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── Helpers ───────────────────────────────────────────────────────────────────
def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank


def fill_bg(slide, color: RGBColor):
    """Fill the slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()          # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.color.rgb = color
    run.font.name    = font_name
    return txb


def add_image_fit(slide, img_path, left, top, width, height):
    """Add image cropped/fitted to the given bounding box."""
    if not os.path.exists(img_path):
        print(f"  ⚠️  Imagem não encontrada: {img_path}")
        return
    slide.shapes.add_picture(img_path, left, top, width, height)


def accent_bar(slide, left, top, width=Inches(0.5), height=Pt(3)):
    """Thin green underline accent."""
    add_rect(slide, left, top, width, height, C_GREEN_L)


def tag_label(slide, text, left, top):
    """Small ALL-CAPS tag above the title."""
    add_text(slide, text.upper(), left, top, Inches(6), Pt(18),
             font_size=9, bold=True, color=C_GREEN_L, align=PP_ALIGN.LEFT)

def bullet_card(slide, icon, title, body, left, top, w=Inches(2.9), h=Inches(1.55)):
    """Mini card with icon + title + body."""
    add_rect(slide, left, top, w, h, C_ACCENT)
    # icon
    add_text(slide, icon, left + Inches(0.18), top + Inches(0.12),
             Inches(0.5), Inches(0.4), font_size=18, color=C_WHITE)
    # title
    add_text(slide, title,
             left + Inches(0.18), top + Inches(0.45),
             w - Inches(0.3), Inches(0.32),
             font_size=11, bold=True, color=C_WHITE)
    # body
    add_text(slide, body,
             left + Inches(0.18), top + Inches(0.8),
             w - Inches(0.3), h - Inches(0.85),
             font_size=9, color=C_MUTED)


def metric_card(slide, icon, value, label, left, top, accent=C_GREEN_L):
    w, h = Inches(3.7), Inches(2.0)
    add_rect(slide, left, top, w, h, C_ACCENT)
    add_rect(slide, left, top + h - Pt(4), w, Pt(4), accent)
    add_text(slide, icon, left + Inches(0.2), top + Inches(0.18),
             Inches(0.6), Inches(0.5), font_size=22, color=C_WHITE)
    add_text(slide, value,
             left + Inches(0.2), top + Inches(0.6),
             w - Inches(0.3), Inches(0.75),
             font_size=32, bold=True, color=accent, font_name="Calibri")
    add_text(slide, label,
             left + Inches(0.2), top + Inches(1.35),
             w - Inches(0.3), Inches(0.55),
             font_size=10, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_intro(prs):
    """Slide 1 – Introdução com logo real, centrado e impactante."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, C_BG_DARK)

    # Full-width top green bar
    add_rect(slide, 0, 0, W, Inches(0.07), C_GREEN)
    # Full-width bottom green bar
    add_rect(slide, 0, H - Inches(0.07), W, Inches(0.07), C_GREEN)

    # Subtle background card in centre
    add_rect(slide, Inches(2.8), Inches(0.6), Inches(7.7), Inches(6.3), C_ACCENT)

    # Logo image – centred
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH,
                                 Inches(5.6), Inches(0.9),
                                 Inches(2.1), Inches(2.1))

    # Brand name
    add_text(slide, "EcoKambio",
             Inches(2.8), Inches(3.1), Inches(7.7), Inches(1.1),
             font_size=72, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, font_name="Calibri")

    # Tagline
    add_text(slide,
             "A Plataforma de Câmbio e Finanças Digitais de Angola",
             Inches(2.8), Inches(4.15), Inches(7.7), Inches(0.55),
             font_size=18, color=C_GREEN_L,
             align=PP_ALIGN.CENTER)

    # Divider line
    add_rect(slide, Inches(5.2), Inches(4.82), Inches(2.9), Pt(2), C_GREEN)

    # Apresentação label
    add_text(slide,
             "APRESENTAÇÃO PARA INVESTIDORES  ·  2026",
             Inches(2.8), Inches(5.0), Inches(7.7), Inches(0.4),
             font_size=11, color=C_MUTED,
             align=PP_ALIGN.CENTER)

    # Website
    add_text(slide, "www.ecokambio.com",
             Inches(2.8), Inches(5.5), Inches(7.7), Inches(0.35),
             font_size=12, color=C_GREEN_L,
             align=PP_ALIGN.CENTER)

    # Confidential stamp bottom-right
    add_text(slide, "CONFIDENCIAL",
             Inches(10.5), Inches(6.9), Inches(2.6), Inches(0.4),
             font_size=9, bold=True, color=C_MUTED, align=PP_ALIGN.RIGHT)


def slide_final(prs):
    """Slide final – Encerramento profissional."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, C_BG_DARK)

    # Full-width top and bottom green bars
    add_rect(slide, 0, 0, W, Inches(0.07), C_GREEN)
    add_rect(slide, 0, H - Inches(0.07), W, Inches(0.07), C_GREEN)

    # Left dark panel
    add_rect(slide, 0, Inches(0.07), Inches(6.5), H - Inches(0.14), C_ACCENT)

    # Logo on left panel – centred vertically
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH,
                                 Inches(2.2), Inches(1.3),
                                 Inches(2.1), Inches(2.1))

    add_text(slide, "EcoKambio",
             Inches(0.5), Inches(3.55), Inches(5.5), Inches(0.9),
             font_size=44, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, font_name="Calibri")

    add_text(slide, "Finanças Digitais · Angola",
             Inches(0.5), Inches(4.38), Inches(5.5), Inches(0.45),
             font_size=13, color=C_GREEN_L, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(1.5), Inches(4.95), Inches(3.5), Pt(2), C_GREEN)

    add_text(slide, "www.ecokambio.com",
             Inches(0.5), Inches(5.15), Inches(5.5), Inches(0.38),
             font_size=13, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Right panel – thank you + contact
    add_text(slide, "Obrigado.",
             Inches(7.0), Inches(1.6), Inches(5.8), Inches(1.3),
             font_size=58, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri")

    add_rect(slide, Inches(7.0), Inches(2.95), Inches(0.55), Pt(3), C_GREEN_L)

    add_text(slide,
             "Agradecemos o seu tempo e interesse\nno EcoKambio.",
             Inches(7.0), Inches(3.1), Inches(5.8), Inches(0.75),
             font_size=14, color=C_MUTED, align=PP_ALIGN.LEFT)

    # Contact boxes
    contacts = [
        ("📩", "Email", "suporte@ecokambio.com"),
        ("🌐", "Website", "www.ecokambio.com"),
        ("💬", "WhatsApp", "Contacto directo via plataforma"),
    ]
    for i, (icon, label, val) in enumerate(contacts):
        cy = Inches(4.1) + i * Inches(0.88)
        add_rect(slide, Inches(7.0), cy, Inches(5.8), Inches(0.75), C_ACCENT)
        add_rect(slide, Inches(7.0), cy, Pt(4), Inches(0.75), C_GREEN_L)
        add_text(slide, f"{icon}  {label}",
                 Inches(7.15), cy + Inches(0.07),
                 Inches(1.3), Inches(0.3),
                 font_size=9, bold=True, color=C_GREEN_L)
        add_text(slide, val,
                 Inches(7.15), cy + Inches(0.37),
                 Inches(5.4), Inches(0.28),
                 font_size=11, color=C_WHITE)

    add_text(slide, "CONFIDENCIAL  ·  2026",
             Inches(7.0), Inches(6.9), Inches(5.8), Inches(0.35),
             font_size=9, color=C_MUTED, align=PP_ALIGN.LEFT)


def slide_screenshot_right(prs, tag, title, subtitle, img_key, cards):
    """Slide template: title left, screenshot right, 3 bullet cards bottom-left."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.06), H, C_GREEN)

    # --- Header area ---
    tag_label(slide, tag, Inches(0.5), Inches(0.32))
    add_text(slide, title,
             Inches(0.5), Inches(0.55), Inches(6), Inches(1.0),
             font_size=32, bold=True, color=C_WHITE, font_name="Calibri")
    accent_bar(slide, Inches(0.5), Inches(1.58), Inches(0.6), Pt(3))
    add_text(slide, subtitle,
             Inches(0.5), Inches(1.7), Inches(5.6), Inches(0.7),
             font_size=11, color=C_MUTED)

    # --- Screenshot (right 55%) ---
    img_path = IMG.get(img_key, "")
    add_image_fit(slide, img_path,
                  Inches(6.3), Inches(0.2),
                  Inches(6.8), Inches(6.95))

    # --- Bullet cards (bottom-left, 3 cols) ---
    card_y = Inches(2.6)
    card_x_start = Inches(0.3)
    gap = Inches(0.12)
    cw  = Inches(1.9)
    for i, (icon, ttl, body) in enumerate(cards[:3]):
        bullet_card(slide, icon, ttl, body,
                    card_x_start + i * (cw + gap), card_y,
                    w=cw, h=Inches(1.95))

    # footer
    add_text(slide, "ecokambio.com", Inches(0.5), Inches(7.1), Inches(4), Pt(16),
             font_size=9, color=C_MUTED)


def slide_metrics(prs):
    """Slide 7 – Métricas & Impacto."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.06), H, C_GREEN)

    tag_label(slide, "Sobre Nós", Inches(0.5), Inches(0.32))
    add_text(slide, "Missão & Impacto",
             Inches(0.5), Inches(0.55), Inches(12), Inches(0.9),
             font_size=36, bold=True, color=C_WHITE, font_name="Calibri")
    accent_bar(slide, Inches(0.5), Inches(1.48), Inches(0.6), Pt(3))
    add_text(slide,
             "Democratizar a informação financeira em Angola — dar a cada angolano acesso "
             "às mesmas ferramentas disponíveis nos mercados desenvolvidos.",
             Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.55),
             font_size=11, color=C_MUTED)

    metrics = [
        ("👥", "50K+",  "Utilizadores únicos mensais",   C_GREEN_L),
        ("🧮", "100K+", "Cálculos de importação feitos",  C_GOLD),
        ("🏦", "6",     "Bancos e fontes monitorizados",  RGBColor(0x41,0x89,0xE0)),
        ("📍", "18",    "Províncias angolanas cobertas",  RGBColor(0x0D,0x94,0x88)),
        ("⏱️", "24/7",  "Monitorização automática",       RGBColor(0x7C,0x3A,0xED)),
        ("🆓", "100%",  "Gratuito para o utilizador",     RGBColor(0xEA,0x58,0x0C)),
    ]

    cols = 3
    mx0 = Inches(0.4)
    my0 = Inches(2.35)
    mgx = Inches(0.22)
    mgy = Inches(0.2)
    mw  = Inches(3.7)
    mh  = Inches(2.0)

    for i, (icon, val, lbl, color) in enumerate(metrics):
        col = i % cols
        row = i // cols
        x = mx0 + col * (mw + mgx)
        y = my0 + row * (mh + mgy)
        metric_card(slide, icon, val, lbl, x, y, accent=color)

    add_text(slide, "ecokambio.com", Inches(0.5), Inches(7.1), Inches(4), Pt(16),
             font_size=9, color=C_MUTED)


def slide_faq_cta(prs):
    """Slide 8 – FAQ + CTA final."""
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, C_BG_DARK)
    add_rect(slide, 0, 0, Inches(0.06), H, C_GREEN)

    tag_label(slide, "Suporte & Transparência", Inches(0.5), Inches(0.32))
    add_text(slide, "FAQ & Próximos Passos",
             Inches(0.5), Inches(0.55), Inches(6), Inches(0.9),
             font_size=32, bold=True, color=C_WHITE, font_name="Calibri")
    accent_bar(slide, Inches(0.5), Inches(1.48), Inches(0.6), Pt(3))
    add_text(slide,
             "Transparência total com os nossos utilizadores — desde como funciona "
             "o cartão Visa até como calculamos os preços de importação.",
             Inches(0.5), Inches(1.6), Inches(5.6), Inches(0.7),
             font_size=11, color=C_MUTED)

    # FAQ screenshot
    add_image_fit(slide, IMG["faq"],
                  Inches(6.3), Inches(0.2),
                  Inches(6.8), Inches(6.95))

    # Próximos passos card
    next_items = [
        "📱  App Mobile Nativa (Android & iOS)",
        "🔔  Alertas de câmbio personalizados",
        "🌍  Mais divisas africanas (ZAR, MZN...)",
        "🏦  Integração bancária directa",
        "💹  Histórico e gráficos de taxas",
    ]
    add_rect(slide, Inches(0.3), Inches(2.5), Inches(5.7), Inches(2.5), C_ACCENT)
    add_rect(slide, Inches(0.3), Inches(2.5), Pt(4), Inches(2.5), C_GREEN_L)
    add_text(slide, "🚀  Próximos Passos",
             Inches(0.5), Inches(2.6), Inches(5.3), Inches(0.4),
             font_size=12, bold=True, color=C_GREEN_L)
    add_text(slide, "\n".join(next_items),
             Inches(0.5), Inches(3.05), Inches(5.3), Inches(1.85),
             font_size=11, color=C_MUTED)

    # Contacto card
    add_rect(slide, Inches(0.3), Inches(5.15), Inches(5.7), Inches(1.3), C_ACCENT)
    add_rect(slide, Inches(0.3), Inches(5.15), Pt(4), Inches(1.3), C_GOLD)
    add_text(slide, "📩  Contacto",
             Inches(0.5), Inches(5.25), Inches(5.3), Inches(0.35),
             font_size=12, bold=True, color=C_GOLD)
    add_text(slide, "suporte@ecokambio.com   ·   www.ecokambio.com",
             Inches(0.5), Inches(5.65), Inches(5.3), Inches(0.65),
             font_size=11, color=C_MUTED)

    add_text(slide, "ecokambio.com", Inches(0.5), Inches(7.1), Inches(4), Pt(16),
             font_size=9, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    print("🟢  A gerar EcoKambio_Investidores_2026.pptx …")
    prs = new_prs()

    # 1 – Introdução com logo
    print("  Slide 1 – Introdução")
    slide_intro(prs)

    # 2 – Mercado Informal
    print("  Slide 2 – Mercado Informal")
    slide_screenshot_right(
        prs,
        tag   = "Mercado de Câmbio",
        title = "Mercado Informal em Tempo Real",
        subtitle = (
            "USD, EUR, USDT, ZAR, BRL e GBP vs AOA — cotações actualizadas automaticamente "
            "via Binance e Coinbase, com conversor integrado."
        ),
        img_key = "mercado_informal",
        cards = [
            ("📊", "6 Divisas",        "USD, EUR, USDT, ZAR, BRL e GBP monitorizadas em simultâneo."),
            ("🔄", "Conversor Rápido", "Converta qualquer valor para AOA instantaneamente, 100% gratuito."),
            ("⏱️", "Sempre Actual",    "Última actualização exibida — total transparência para o utilizador."),
        ]
    )

    # 3 – Produtos tab na homepage
    print("  Slide 3 – Homepage aba Produtos")
    slide_screenshot_right(
        prs,
        tag   = "Marketplace Integrado",
        title = "Produtos & Serviços na Homepage",
        subtitle = (
            "Visa Virtual, EcoFlix e centenas de produtos importados acessíveis "
            "directamente na página principal — sem sair do ecossistema."
        ),
        img_key = "homepage_produtos",
        cards = [
            ("💳", "Visa Virtual",   "Cartão digital para compras globais, entregue instantaneamente."),
            ("🎬", "EcoFlix",        "Netflix Angola com planos partilhados a preços acessíveis."),
            ("📦", "Produtos",       "Electrónica, acessórios e mais — preços reais em Kwanzas."),
        ]
    )

    # 4 – Calculadora
    print("  Slide 4 – Calculadora")
    slide_screenshot_right(
        prs,
        tag   = "Ferramenta Inteligente",
        title = "Calculadora de Câmbio",
        subtitle = (
            "Calcula o custo total real da sua compra em Angola — câmbio actualizado, "
            "comissão bancária e taxa de serviço incluídos no resultado."
        ),
        img_key = "calculadora",
        cards = [
            ("🧮", "Resultado Imediato", "109 362 AOA de custo vs 111 855 AOA total: tudo detalhado."),
            ("🏦", "Por Banco",          "Escolha o banco e compare qual oferece a melhor cotação."),
            ("🌍", "Multi-plataforma",   "Para compras no Alibaba, Amazon, AliExpress e muito mais."),
        ]
    )

    # 5 – Parceiros + Produtos
    print("  Slide 5 – Parceiros & Ecossistema")
    slide_screenshot_right(
        prs,
        tag   = "Ecossistema & Parcerias",
        title = "Parceiros & Ecossistema",
        subtitle = (
            "Integração com parceiros de referência como a Bybit — criando um ecossistema "
            "financeiro completo para o utilizador angolano."
        ),
        img_key = "parceiros_produtos",
        cards = [
            ("🤝", "Bybit",          "Parceiro de trading de cripto — acesso a futuros e spot com alta liquidez."),
            ("📦", "Produtos AOA",   "Catálogo de produtos com preços em AOA incluindo todos os custos."),
            ("🔁", "Receita Afiliada","Comissões por cada clique e conversão nos parceiros integrados."),
        ]
    )

    # 6 – EcoFlix
    print("  Slide 6 – EcoFlix")
    slide_screenshot_right(
        prs,
        tag   = "Serviço Digital",
        title = "EcoFlix — Netflix Angola",
        subtitle = (
            "Planos de Netflix partilhados com preços acessíveis em Kwanzas: "
            "Económico (4.500 Kz), Ultra (6.500 Kz) e Família (18.000 Kz)."
        ),
        img_key = "ecoflix",
        cards = [
            ("🎬", "3 Planos",        "Económico, Ultra e Família — para todos os perfis e orçamentos."),
            ("📺", "4K Ultra HDR",    "Qualidade máxima em todos os dispositivos, incluindo Smart TV 4K."),
            ("💰", "Preços em Kz",    "Sem necessidade de cartão estrangeiro — pague directamente em AOA."),
        ]
    )

    # 7 – Visa Virtual
    print("  Slide 7 – Visa Virtual")
    slide_screenshot_right(
        prs,
        tag   = "Serviço Premium",
        title = "Cartão Visa Virtual",
        subtitle = (
            "A porta de entrada para o mundo digital. Compras online e pagamentos "
            "em qualquer lugar do mundo, sem conta bancária."
        ),
        img_key = "visa_hero",
        cards = [
            ("💳", "100% Virtual",    "Número, validade e CVV entregues via WhatsApp em minutos."),
            ("🔒", "Pré-pago Seguro", "Carregue o valor que precisa. 100 USD ≈ 99.750 AOA (taxa 5%)."),
            ("🌐", "Aceite em Tudo",  "Amazon, Netflix, Spotify, Facebook Ads, TikTok Ads e qualquer site Visa."),
        ]
    )

    # 8 – Produto Detalhe
    print("  Slide 8 – Detalhe de Produto")
    slide_screenshot_right(
        prs,
        tag   = "Experiência de Produto",
        title = "Detalhe do Produto",
        subtitle = (
            "Página de produto premium: galeria de imagens, variantes de cor, "
            "preço com desconto em Kz e recomendações inteligentes."
        ),
        img_key = "produto_detalhe",
        cards = [
            ("🎨", "Variantes",       "Cor e configuração com preview visual — Preto, Branco, Verde."),
            ("📸", "Galeria",          "Múltiplas imagens do produto para melhor decisão de compra."),
            ("🔗", "Recomendações",   "Produtos relacionados aumentam o valor médio por sessão."),
        ]
    )

    # 9 – Métricas
    print("  Slide 9 – Missão & Impacto")
    slide_metrics(prs)

    # 10 – Slide Final
    print("  Slide 10 – Encerramento")
    slide_final(prs)

    prs.save(OUT_FILE)
    print(f"\n✅  Ficheiro gerado com sucesso:\n   {OUT_FILE}")


if __name__ == "__main__":
    main()
